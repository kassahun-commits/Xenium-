#!/usr/bin/env python3
"""
Build a digestible PowerPoint overview of the nucleus-only Xenium analysis:
workflow (from the files you hand over -> data), QC parameters, cell typing,
and a Wilcoxon-vs-pseudobulk DE comparison slide.

Generates a couple of purpose-built PNG figures from the DE results CSVs, then
assembles the .pptx. No hardcoded paths (CLI args).
"""
import argparse
from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette -------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x2D, 0x3D)
BLUE = RGBColor(0x1F, 0x77, 0xB4)
RED = RGBColor(0xD6, 0x27, 0x28)
GRAY = RGBColor(0x55, 0x55, 0x55)
LGRAY = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
AMBER = RGBColor(0xE8, 0x8A, 0x00)

GROUP_COLORS = {
    'H2O_veh': '#7f7f7f', 'H2O_MCT1i': '#17becf', 'EtOH_veh': '#ff7f0e',
    'EtOH_MCT1i': '#2ca02c', 'ChronicEtOH': '#d62728',
    'MAT2A_CM': '#9467bd', 'MAT2A_OE': '#8c564b',
}
COMP_ORDER = ['H2O_MCT1i', 'EtOH_veh', 'EtOH_MCT1i', 'ChronicEtOH',
              'MAT2A_CM', 'MAT2A_OE']


# ---- figure helpers ------------------------------------------------------
def fig_sig_bar(pb_csv, wx_csv, out_png):
    """Grouped bar: # significant genes per comparison, Wilcoxon vs pseudobulk,
    one panel per cell type."""
    pb = pd.read_csv(pb_csv)
    wx = pd.read_csv(wx_csv)

    def counts(df):
        d = df[df['significant']].groupby(['celltype', 'comparison']).size()
        return d

    cpb, cwx = counts(pb), counts(wx)
    cts = ['Neuron', 'Astrocyte']
    fig, axes = plt.subplots(1, 2, figsize=(10, 3.6), sharey=False)
    for ax, ct in zip(axes, cts):
        labels = COMP_ORDER
        pb_vals = [int(cpb.get((ct, f'{c}_vs_H2O_veh'), 0)) for c in labels]
        wx_vals = [int(cwx.get((ct, f'{c}_vs_H2O_veh'), 0)) for c in labels]
        x = np.arange(len(labels))
        w = 0.38
        b1 = ax.bar(x - w / 2, wx_vals, w, color='#d62728', label='Cell-level Wilcoxon')
        b2 = ax.bar(x + w / 2, pb_vals, w, color='#1f77b4', label='Pseudobulk DESeq2')
        ax.set_title(ct, fontsize=12, fontweight='bold')
        ax.set_xticks(x)
        ax.set_xticklabels([c.replace('_', '\n') for c in labels], fontsize=7)
        ax.set_ylabel('# significant genes', fontsize=9)
        ax.bar_label(b1, fontsize=7, padding=1)
        ax.bar_label(b2, fontsize=7, padding=1)
        for s in ['top', 'right']:
            ax.spines[s].set_visible(False)
        if ct == 'Neuron':
            ax.legend(fontsize=8, frameon=False, loc='upper left')
    fig.suptitle('Significant genes vs H2O_veh  (|log2FC|>0.5)', fontsize=12)
    plt.tight_layout(rect=[0, 0, 1, 0.95])
    plt.savefig(out_png, dpi=170, bbox_inches='tight')
    plt.close(fig)


def fig_volcano_pair(pb_csv, wx_csv, celltype, treatment, out_png):
    """Side-by-side volcano (pseudobulk | Wilcoxon) for one comparison."""
    comp = f'{treatment}_vs_H2O_veh'
    pb = pd.read_csv(pb_csv)
    wx = pd.read_csv(wx_csv)
    fig, axes = plt.subplots(1, 2, figsize=(10, 4), sharex=True)
    for ax, df, name, padj_cut in [
        (axes[0], pb, 'Pseudobulk DESeq2\n(padj<0.05)', 0.05),
        (axes[1], wx, 'Cell-level Wilcoxon\n(padj<0.01) — EXPLORATORY', 0.01)]:
        d = df[(df['celltype'] == celltype) & (df['comparison'] == comp)].copy()
        d = d.dropna(subset=['log2FoldChange', 'padj'])
        x = d['log2FoldChange'].values
        y = -np.log10(np.clip(d['padj'].values, 1e-300, 1))
        sig = (d['padj'] < padj_cut) & (d['log2FoldChange'].abs() > 0.5)
        ax.scatter(x[~sig.values], y[~sig.values], s=5, c='#cccccc',
                   alpha=0.5, edgecolor='none', rasterized=True)
        ax.scatter(x[sig.values], y[sig.values], s=10,
                   c='#d62728' if 'Wilcoxon' in name else '#1f77b4',
                   alpha=0.8, edgecolor='none')
        ax.axhline(-np.log10(padj_cut), ls='--', color='k', lw=0.5, alpha=0.4)
        ax.axvline(0.5, ls='--', color='k', lw=0.5, alpha=0.4)
        ax.axvline(-0.5, ls='--', color='k', lw=0.5, alpha=0.4)
        ax.set_title(f'{name}\n{int(sig.sum())} significant', fontsize=10)
        ax.set_xlabel('log2 fold change', fontsize=9)
        ax.set_ylabel('-log10 adj. p', fontsize=9)
        for s in ['top', 'right']:
            ax.spines[s].set_visible(False)
    fig.suptitle(f'{celltype}: {treatment} vs H2O_veh — same data, two tests',
                 fontsize=12, fontweight='bold')
    plt.tight_layout(rect=[0, 0, 1, 0.93])
    plt.savefig(out_png, dpi=170, bbox_inches='tight')
    plt.close(fig)


# ---- pptx helpers --------------------------------------------------------
def add_title_only(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def txt(slide, l, t, w, h, text, size=18, bold=False, color=NAVY,
        align=PP_ALIGN.LEFT, font='Calibri'):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = font
    return tb


def bullets(slide, l, t, w, h, items, size=15, color=NAVY, space=6):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (lvl, text, *style) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = ('• ' if lvl == 0 else '– ') + text
        r.font.size = Pt(size - lvl * 1)
        r.font.color.rgb = style[0] if style else color
        r.font.name = 'Calibri'
        if len(style) > 1 and style[1]:
            r.font.bold = True
    return tb


def band(slide, prs, color=NAVY, h=1.05):
    s = slide.shapes.add_shape(1, Inches(0), Inches(0),
                               prs.slide_width, Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def chip(slide, l, t, w, h, text, fill, tcolor=WHITE, size=12, bold=True):
    s = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = WHITE; s.line.width = Pt(1)
    s.shadow.inherit = False
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = tcolor
    r.font.name = 'Calibri'
    return s


def arrow(slide, l, t, w, h=0.3, color=LGRAY):
    s = slide.shapes.add_shape(13, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s


def slide_header(slide, prs, title, sub=None):
    band(slide, prs)
    txt(slide, 0.5, 0.18, 12.3, 0.6, title, size=26, bold=True, color=WHITE)
    if sub:
        txt(slide, 0.5, 0.72, 12.3, 0.35, sub, size=13, color=RGBColor(0xCF, 0xDA, 0xE6))


# ---- build ---------------------------------------------------------------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--nuc-dir', required=True, help='nucleus_only output folder')
    ap.add_argument('--date', default='2026-06-09')
    ap.add_argument('--out', default=None)
    args = ap.parse_args()

    nd = Path(args.nuc_dir)
    D = args.date
    assets = nd / 'deck_assets'; assets.mkdir(exist_ok=True)
    pb_csv = nd / f'Nucleus_PseudobulkDE_results_{D}.csv'
    wx_csv = nd / f'Nucleus_WilcoxonDE_results_{D}.csv'
    qc_pdf = nd / f'SlidesAB_nucleus_QCverification_{D}.pdf'

    # figures
    sig_png = assets / 'sig_bar.png'
    pair_png = assets / 'volcano_pair_Neuron_EtOH_veh.png'
    fig_sig_bar(pb_csv, wx_csv, sig_png)
    fig_volcano_pair(pb_csv, wx_csv, 'Neuron', 'EtOH_veh', pair_png)
    qc_png = assets / 'qc.png'
    import subprocess
    subprocess.run(['sips', '-s', 'format', 'png', str(qc_pdf),
                    '--out', str(qc_png)], capture_output=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: title ----
    s = add_title_only(prs)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
    bg.shadow.inherit = False
    txt(s, 0.8, 2.2, 11.7, 1.2, 'Nucleus-only Xenium analysis',
        size=40, bold=True, color=WHITE)
    txt(s, 0.8, 3.35, 11.7, 0.8,
        'From raw output files to differential expression — Slides A + B',
        size=20, color=RGBColor(0xCF, 0xDA, 0xE6))
    txt(s, 0.8, 6.4, 11.7, 0.5,
        f'Project: 10X nucleus only   ·   {D}   ·   kassahun@bu.edu',
        size=13, color=LGRAY)

    # ---- Slide 2: workflow from files -> data ----
    s = add_title_only(prs)
    slide_header(s, prs, 'What happens after you hand me the files',
                 'The Xenium Explorer CSVs are only ROI outlines — the data itself is rebuilt from the output folder')
    # two source boxes
    chip(s, 0.5, 1.45, 3.1, 1.0,
         'Xenium Explorer CSVs\n(Slide_A / Slide_B annotations)\n→ ROI polygon coordinates ONLY',
         RGBColor(0x6B, 0x4E, 0x9E), size=11)
    chip(s, 0.5, 2.7, 3.1, 1.5,
         'Output folder files\n• transcripts.parquet (every molecule:\n  gene, x/y, cell, QV, in-nucleus?)\n• cells.csv.gz (cell centroids)\n• cell_feature_matrix.h5 (gene panel)',
         RGBColor(0x2E, 0x5A, 0x88), size=10.5)
    # pipeline steps
    steps = [
        ('1. Keep nuclear\nmolecules', 'overlaps_nucleus=1,\nQV≥20, real gene,\nassigned to a cell'),
        ('2. Re-count\nper cell × gene', 'build matrix from\ntranscripts (not the\nstandard whole-cell h5)'),
        ('3. Assign ROI', 'point-in-polygon of\ncentroid vs your\nExplorer outlines'),
        ('4. Combine A+B\n+ QC + normalize', 'filter, log-normalize,\none AnnData object'),
        ('5. Cell typing\n+ DE / volcanoes', 'marker scores →\ncell type; DESeq2 /\nWilcoxon'),
    ]
    x = 4.1; y = 1.7; bw = 1.72; gap = 0.12
    for i, (head, body) in enumerate(steps):
        chip(s, x, y, bw, 1.0, head, BLUE, size=11)
        txt(s, x, y + 1.05, bw, 1.2, body, size=9.5, color=GRAY)
        if i < len(steps) - 1:
            arrow(s, x + bw + 0.005, y + 0.4, gap + 0.02)
        x += bw + gap
    txt(s, 0.5, 5.7, 12.3, 1.2,
        'Key point: Xenium does not output "the data" as a table of genes per cell type. '
        'It outputs molecule positions + a count matrix. Everything downstream — nucleus '
        'restriction, ROIs, cell types, expression per group — is computed from those files.',
        size=14, color=NAVY)

    # ---- Slide 3: why nucleus rebuild ----
    s = add_title_only(prs)
    slide_header(s, prs, 'Why "nucleus-only" needs a rebuilt matrix')
    bullets(s, 0.6, 1.4, 12.1, 4.5, [
        (0, 'The standard cell_feature_matrix.h5 counts ALL transcripts in a cell — nucleus + cytoplasm.', NAVY, True),
        (0, 'overlaps_nucleus is a per-MOLECULE flag, so "nuclear only" cannot be filtered from the finished matrix.', NAVY),
        (0, 'So we go back to transcripts.parquet (635M molecules across both slides) and re-count, keeping only molecules where:', NAVY),
        (1, 'overlaps_nucleus = 1   (inside the DAPI-segmented nucleus)', GRAY),
        (1, 'QV ≥ 20   (high-confidence calls)', GRAY),
        (1, 'is_gene = True   (drops negative-control / unassigned codewords)', GRAY),
        (1, 'cell_id ≠ UNASSIGNED   (molecule is assigned to a real cell)', GRAY),
        (0, 'Result: 59.7M (Slide A) + 75.0M (Slide B) nuclear molecules → a cell × gene matrix directly comparable to the whole-cell one (same 5,104-gene panel).', NAVY, True),
    ], size=15)

    # ---- Slide 4: QC ----
    s = add_title_only(prs)
    slide_header(s, prs, 'Quality control — what we filtered and why')
    # left: param table
    rows = [
        ('Parameter', 'Value', 'Why'),
        ('Nucleus flag', 'overlaps_nucleus = 1', 'defines "nuclear"'),
        ('QV cutoff', '≥ 20', 'trough of bimodal QV;\n= 10x default'),
        ('Gene panel', '5,104 genes', 'locked to whole-cell\nmatrix (comparable)'),
        ('Min counts/cell', '10', 'matches existing pipeline'),
        ('Min cells/gene', '5', 'matches existing pipeline'),
        ('Normalization', 'CP10k + log1p', 'matches existing pipeline'),
    ]
    tbl_w = 6.2
    gt = s.shapes.add_table(len(rows), 3, Inches(0.5), Inches(1.4),
                            Inches(tbl_w), Inches(3.6)).table
    gt.columns[0].width = Inches(2.0)
    gt.columns[1].width = Inches(1.9)
    gt.columns[2].width = Inches(2.3)
    for j, cell in enumerate(gt.rows[0].cells):
        cell.text = rows[0][j]
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(12)
    for i in range(1, len(rows)):
        for j in range(3):
            c = gt.cell(i, j); c.text = rows[i][j]
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10.5)
                    r.font.color.rgb = NAVY if j < 2 else GRAY
                    if j == 1:
                        r.font.bold = True
    # right: QC figure + result line
    if qc_png.exists():
        s.shapes.add_picture(str(qc_png), Inches(7.0), Inches(1.5), width=Inches(5.9))
    txt(s, 0.5, 5.25, 12.3, 1.6,
        'Result: combined post-QC = 167,428 cells × 5,104 genes (only 130 cells, 0 genes dropped). '
        'Median 533 nuclear transcripts/cell; 18 samples across 7 groups. '
        'QV≥20 was chosen from the data (bimodal distribution, natural valley at 18–20) and because '
        'it matches the cutoff 10x used for the whole-cell matrix — keeping the two directly comparable.',
        size=12.5, color=NAVY)

    # ---- Slide 5: cell typing ----
    s = add_title_only(prs)
    slide_header(s, prs, 'How each cell gets a cell type',
                 'Xenium gives no cell-type labels — they are computed from marker genes')
    bullets(s, 0.6, 1.45, 6.4, 4.8, [
        (0, 'For each lineage, score every cell on a curated marker set (scanpy score_genes).', NAVY, True),
        (0, 'Assign each cell to its highest-scoring lineage (argmax); score ≤ 0 → Unclassified.', NAVY),
        (0, 'Marker sets used:', NAVY),
        (1, 'Neuron: Rbfox3, Snap25, Syn1, Syt1, Stmn2, Map2, Tubb3', BLUE),
        (1, 'Astrocyte: Gfap, Aqp4, Slc1a3, Aldh1l1, S100b, Aldoc, Gja1', RED),
        (1, 'Oligodendrocyte: Mog, Olig1/2, Sox10', RGBColor(0x94,0x67,0xBD)),
        (1, 'Microglia: Cx3cr1, Tmem119, Csf1r, Aif1, Trem2', GREEN),
    ], size=14)
    # counts mini-table
    rows = [('Cell type', 'cells'),
            ('Neuron', '80,940'), ('Astrocyte', '42,437'),
            ('Oligodendrocyte', '25,206'), ('Microglia', '8,204'),
            ('Unclassified', '10,641')]
    gt = s.shapes.add_table(len(rows), 2, Inches(7.4), Inches(1.55),
                            Inches(5.3), Inches(3.4)).table
    for j in range(2):
        c = gt.rows[0].cells[j]; c.text = rows[0][j]
        c.fill.solid(); c.fill.fore_color.rgb = NAVY
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(13)
    for i in range(1, len(rows)):
        for j in range(2):
            c = gt.cell(i, j); c.text = rows[i][j]
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12); r.font.color.rgb = NAVY
                    if j == 0: r.font.bold = True
    txt(s, 7.4, 5.2, 5.3, 1.0,
        'Neurons + astrocytes (the two requested for DE) dominate, so both have\nplenty of cells for analysis.',
        size=11.5, color=GRAY)

    # ---- Slide 6: Wilcoxon vs pseudobulk ----
    s = add_title_only(prs)
    slide_header(s, prs, 'Two ways to test DE: Wilcoxon vs pseudobulk',
                 'Same cells, same thresholds on fold-change — very different significance')
    s.shapes.add_picture(str(sig_png), Inches(0.4), Inches(1.25), width=Inches(8.4))
    bullets(s, 9.0, 1.4, 4.1, 4.4, [
        (0, 'Cell-level Wilcoxon treats each cell as a sample → thinks n = thousands.', RED, True),
        (0, 'With only 2–3 mice/group that is pseudoreplication: p-values track cell counts, not biology → inflated hit lists.', NAVY),
        (0, 'Pseudobulk DESeq2 sums counts per mouse and tests across replicates (n=2–3) → the rigorous answer.', BLUE, True),
        (0, 'Wilcoxon = continuity with the prior grant figures (EXPLORATORY).', GRAY),
        (0, 'Caveat for BOTH: MAT2A groups are Slide-B-only → those hits are confounded with slide/batch.', AMBER, True),
    ], size=12.5)

    # ---- Slide 7: example volcano pair ----
    s = add_title_only(prs)
    slide_header(s, prs, 'Same comparison, two methods',
                 'Neuron — EtOH_veh vs H2O_veh: 559 "hits" by Wilcoxon vs 0 by pseudobulk')
    s.shapes.add_picture(str(pair_png), Inches(1.4), Inches(1.5), width=Inches(10.5))
    txt(s, 0.6, 6.55, 12.1, 0.8,
        'The 559 Wilcoxon "significant" genes are an artifact of counting cells as replicates. '
        'Across mice, the effect is not reproducible — so pseudobulk flags none. This is why pseudobulk is the primary analysis.',
        size=12.5, color=NAVY)

    # ---- Slide 8: outputs / takeaways ----
    s = add_title_only(prs)
    slide_header(s, prs, 'Takeaways & where everything lives')
    bullets(s, 0.6, 1.4, 12.1, 3.2, [
        (0, 'Nucleus-only object: SlidesAB_nucleus_combined_2026-06-09.h5ad (167,428 cells × 5,104 genes).', NAVY, True),
        (0, 'Two DE result sets + volcanoes (Neuron, Astrocyte; all groups vs H2O_veh): pseudobulk (primary) and Wilcoxon (grant continuity).', NAVY),
        (0, 'Every figure has editable text + a matching source-data CSV; scripts are on GitHub (kassahun-commits/Xenium-).', NAVY),
        (0, 'For the grant: lead with pseudobulk; keep Wilcoxon for continuity; state n=2–3 as pilot data and note the MAT2A slide confound.', BLUE, True),
    ], size=14)
    txt(s, 0.6, 5.0, 12.1, 1.8,
        'Folder: analysis/current/nucleus_only/\n'
        'Scripts: build_nucleus_matrix_AB.py · nucleus_pseudobulk_volcano_AB.py · nucleus_wilcoxon_volcano_AB.py',
        size=12, color=GRAY)

    out = Path(args.out) if args.out else nd / f'Nucleus_only_workflow_overview_{D}.pptx'
    prs.save(str(out))
    print(f'[saved] {out}')
    print(f'[assets] {sig_png.name}, {pair_png.name}, {qc_png.name}')


if __name__ == '__main__':
    main()

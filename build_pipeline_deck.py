#!/usr/bin/env python3
"""
Build a GENERAL teaching PowerPoint on running a 10x Xenium analysis with Claude,
plus the statistical caveats to watch for. Broad audience, large fonts.

Sections (confirmed with user 2026-06-09):
  1. The workflow, from the CSV/output files to results (general, large fonts).
  2-3. QC parameters to have Claude check: what each means + recommended threshold;
       plus how to filter by cluster in Xenium Explorer (light steps).
  3b.  Cell typing: turning anonymous clusters into named cell types via marker
       genes (recipe + brain marker cheat-sheet + where the lists come from).
  3c.  Spatial neighborhood / niche analysis.
  4. Wilcoxon vs pseudobulk: why/when each, + the "null pseudobulk != no effect"
     diagnostic (effects go the same direction, just underpowered).
  5. The Wilcoxon-vs-pseudobulk significant-gene bar chart.

Regenerates two figures from the V1 whole-cell DE result CSVs, then assembles the
deck. No hardcoded paths (CLI args).
"""
import argparse
import subprocess
from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

plt.rcParams['pdf.fonttype'] = 42
plt.rcParams['ps.fonttype'] = 42

# ---- palette -------------------------------------------------------------
NAVY = RGBColor(0x1F, 0x2D, 0x3D)
BLUE = RGBColor(0x1F, 0x77, 0xB4)
RED = RGBColor(0xD6, 0x27, 0x28)
GRAY = RGBColor(0x55, 0x55, 0x55)
LGRAY = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
AMBER = RGBColor(0xE8, 0x8A, 0x00)
PURPLE = RGBColor(0x6B, 0x4E, 0x9E)
TEAL = RGBColor(0x0E, 0x7C, 0x86)   # tag: whole-cell segmentation
NUCP = RGBColor(0x8E, 0x44, 0xAD)   # tag: nucleus-only segmentation
SLATE = RGBColor(0x44, 0x55, 0x66)  # tag: "choice / axis" explainer slides

# V1 whole-cell contrasts (test, reference) -> pseudobulk comparison string
V1_CONTRASTS = [
    ('EtOH_veh', 'H2O_veh'),
    ('ChronicEtOH', 'H2O_veh'),
    ('H2O_MCT1i', 'H2O_veh'),
    ('EtOH_MCT1i', 'H2O_veh'),
    ('EtOH_MCT1i', 'EtOH_veh'),
    ('MAT2A_OE', 'MAT2A_CM'),
    ('APP', 'H2O_veh'),
]
# nucleus-only contrasts (all vs H2O_veh) for the nucleus bar slide
NUC_CONTRASTS = [
    ('EtOH_veh', 'H2O_veh'),
    ('ChronicEtOH', 'H2O_veh'),
    ('H2O_MCT1i', 'H2O_veh'),
    ('EtOH_MCT1i', 'H2O_veh'),
    ('MAT2A_CM', 'H2O_veh'),
    ('MAT2A_OE', 'H2O_veh'),
]
PB_PADJ, PB_LFC = 0.05, 0.5
WX_PADJ, WX_LFC = 1e-3, 0.5
DIAG_CT, DIAG_TEST, DIAG_REF = 'Neuron', 'EtOH_veh', 'H2O_veh'


# ---- figure helpers ------------------------------------------------------
def _wx_sig(wx):
    return (wx['padj'] < WX_PADJ) & (wx['logfc'].abs() > WX_LFC)


def fig_sig_bar(pb_csv, wx_csv, out_png):
    """Grouped bar: # significant genes per V1 contrast, Wilcoxon vs pseudobulk,
    one panel per cell type."""
    pb = pd.read_csv(pb_csv)
    wx = pd.read_csv(wx_csv)
    wx = wx.copy()
    wx['sig'] = _wx_sig(wx)

    cts = ['Neuron', 'Astrocyte']
    labels = [f'{t}\nvs {r}' for t, r in V1_CONTRASTS]
    fig, axes = plt.subplots(1, 2, figsize=(11.5, 4.0))
    for ax, ct in zip(axes, cts):
        pb_vals, wx_vals = [], []
        for t, r in V1_CONTRASTS:
            comp = f'{t}_vs_{r}'
            pb_vals.append(int(pb[(pb['celltype'] == ct) &
                                  (pb['comparison'] == comp) &
                                  (pb['significant'])].shape[0]))
            wx_vals.append(int(wx[(wx['subtype'] == ct) & (wx['test'] == t) &
                                  (wx['reference'] == r) & (wx['sig'])].shape[0]))
        x = np.arange(len(labels))
        w = 0.4
        b1 = ax.bar(x - w / 2, wx_vals, w, color='#d62728',
                    label='Cell-level Wilcoxon (padj<1e-3)')
        b2 = ax.bar(x + w / 2, pb_vals, w, color='#1f77b4',
                    label='Pseudobulk DESeq2 (padj<0.05)')
        ax.set_title(ct, fontsize=14, fontweight='bold')
        ax.set_xticks(x)
        ax.set_xticklabels(labels, fontsize=7.5)
        ax.set_ylabel('# significant genes (|log2FC|>0.5)', fontsize=10)
        ax.bar_label(b1, fontsize=7.5, padding=1)
        ax.bar_label(b2, fontsize=7.5, padding=1)
        for sp in ['top', 'right']:
            ax.spines[sp].set_visible(False)
        if ct == 'Neuron':
            ax.legend(fontsize=9, frameon=False, loc='upper left')
    fig.suptitle('Same cells, two tests — Wilcoxon inflates the hit list',
                 fontsize=14, fontweight='bold')
    plt.tight_layout(rect=[0, 0, 1, 0.94])
    plt.savefig(out_png, dpi=170, bbox_inches='tight')
    plt.close(fig)


def fig_diagnostic(pb_csv, wx_csv, out_png):
    """Two panels for the EtOH_veh vs H2O_veh neuron contrast:
    (L) pseudobulk volcano showing big-but-underpowered effects below the sig line;
    (R) the diagnostic metrics that say 'underpowered, not null'."""
    pb = pd.read_csv(pb_csv)
    wx = pd.read_csv(wx_csv)
    comp = f'{DIAG_TEST}_vs_{DIAG_REF}'
    d = pb[(pb['celltype'] == DIAG_CT) & (pb['comparison'] == comp)].copy()
    d = d.dropna(subset=['log2FoldChange', 'padj'])

    # metrics
    n_lfc05 = int((d['log2FoldChange'].abs() > 0.5).sum())
    n_lfc1 = int((d['log2FoldChange'].abs() > 1.0).sum())
    min_padj = float(d['padj'].min())
    w = wx[(wx['subtype'] == DIAG_CT) & (wx['test'] == DIAG_TEST) &
           (wx['reference'] == DIAG_REF)].copy()
    wsig = w[_wx_sig(w)]
    m = wsig.merge(d[['gene', 'log2FoldChange']], on='gene')
    conc = (np.sign(m['logfc']) == np.sign(m['log2FoldChange'])).mean() * 100 \
        if len(m) else float('nan')

    fig, axes = plt.subplots(1, 2, figsize=(11.5, 4.4),
                             gridspec_kw={'width_ratios': [1.25, 1]})
    # ---- L: volcano ----
    ax = axes[0]
    x = d['log2FoldChange'].values
    y = -np.log10(np.clip(d['padj'].values, 1e-300, 1))
    sig = (d['padj'] < PB_PADJ) & (d['log2FoldChange'].abs() > PB_LFC)
    ax.scatter(x[~sig.values], y[~sig.values], s=7, c='#cccccc', alpha=0.55,
               edgecolor='none', rasterized=True)
    ax.scatter(x[sig.values], y[sig.values], s=12, c='#1f77b4',
               edgecolor='none')
    ax.axhline(-np.log10(PB_PADJ), ls='--', color='k', lw=0.7, alpha=0.5)
    ax.text(ax.get_xlim()[1], -np.log10(PB_PADJ), ' padj=0.05',
            fontsize=8, va='bottom', ha='right', color='k')
    for v in (-PB_LFC, PB_LFC):
        ax.axvline(v, ls='--', color='k', lw=0.7, alpha=0.5)
    # label the largest-effect near-miss genes
    cand = d[d['log2FoldChange'].abs() > 1.0].nsmallest(8, 'padj')
    for _, rr in cand.iterrows():
        ax.annotate(rr['gene'],
                    (rr['log2FoldChange'], -np.log10(max(rr['padj'], 1e-300))),
                    fontsize=8, ha='center', va='bottom', color='#444')
    ax.set_title(f'{DIAG_CT}: {DIAG_TEST} vs {DIAG_REF}\n'
                 f'pseudobulk — 0 pass, but effects are real-sized',
                 fontsize=11)
    ax.set_xlabel('log2 fold change', fontsize=10)
    ax.set_ylabel('-log10 adjusted p', fontsize=10)
    for sp in ['top', 'right']:
        ax.spines[sp].set_visible(False)

    # ---- R: metrics ----
    ax = axes[1]
    ax.axis('off')
    lines = [
        ('Is the null real, or underpowered?', '#1F2D3D', 13, True),
        ('', '#1F2D3D', 6, False),
        (f'• {n_lfc05:,} genes with |log2FC| > 0.5', '#333333', 12, False),
        (f'• {n_lfc1:,} genes with |log2FC| > 1.0  (large!)', '#333333', 12, False),
        (f'• best gene only reaches padj = {min_padj:.2f}', '#333333', 12, False),
        ('   (not ~1.0 — knocking on the door)', '#777777', 10, False),
        (f'• {conc:.0f}% of Wilcoxon hits go the SAME', '#1f77b4', 12, True),
        ('   direction in pseudobulk', '#1f77b4', 12, True),
        ('', '#1F2D3D', 8, False),
        ('Verdict: effects are present and consistent,', '#2ca02c', 12, True),
        ('but n=2 vs 2 is too few to certify them.', '#2ca02c', 12, True),
        ('→ A null here means ADD Ns, not "no effect".', '#d62728', 12, True),
    ]
    yy = 0.97
    for text, col, sz, bold in lines:
        if text:
            ax.text(0.0, yy, text, fontsize=sz, color=col,
                    fontweight='bold' if bold else 'normal',
                    transform=ax.transAxes, va='top')
        yy -= 0.075 if text else 0.04
    fig.suptitle('Reading a "no significant DE" pseudobulk result',
                 fontsize=14, fontweight='bold')
    plt.tight_layout(rect=[0, 0, 1, 0.93])
    plt.savefig(out_png, dpi=170, bbox_inches='tight')
    plt.close(fig)


def fig_sig_bar_nucleus(pb_csv, wx_csv, out_png):
    """Nucleus-only version of the V1 sig bar: # significant genes per contrast,
    Wilcoxon vs pseudobulk, one panel per cell type. Both nucleus CSVs share a
    `comparison` (= '{test}_vs_{ref}') schema and a `log2FoldChange` column."""
    pb = pd.read_csv(pb_csv)
    wx = pd.read_csv(wx_csv)
    pb['sig'] = (pb['padj'] < PB_PADJ) & (pb['log2FoldChange'].abs() > PB_LFC)
    wx['sig'] = (wx['padj'] < WX_PADJ) & (wx['log2FoldChange'].abs() > WX_LFC)

    cts = ['Neuron', 'Astrocyte']
    labels = [f'{t}\nvs {r}' for t, r in NUC_CONTRASTS]
    fig, axes = plt.subplots(1, 2, figsize=(11.5, 4.0))
    for ax, ct in zip(axes, cts):
        pb_vals, wx_vals = [], []
        for t, r in NUC_CONTRASTS:
            comp = f'{t}_vs_{r}'
            pb_vals.append(int(pb[(pb['celltype'] == ct) &
                                  (pb['comparison'] == comp) &
                                  (pb['sig'])].shape[0]))
            wx_vals.append(int(wx[(wx['celltype'] == ct) &
                                  (wx['comparison'] == comp) &
                                  (wx['sig'])].shape[0]))
        x = np.arange(len(labels))
        w = 0.4
        b1 = ax.bar(x - w / 2, wx_vals, w, color='#d62728',
                    label='Cell-level Wilcoxon (padj<1e-3)')
        b2 = ax.bar(x + w / 2, pb_vals, w, color='#1f77b4',
                    label='Pseudobulk DESeq2 (padj<0.05)')
        ax.set_title(ct, fontsize=14, fontweight='bold')
        ax.set_xticks(x)
        ax.set_xticklabels(labels, fontsize=7.5)
        ax.set_ylabel('# significant genes (|log2FC|>0.5)', fontsize=10)
        ax.bar_label(b1, fontsize=7.5, padding=1)
        ax.bar_label(b2, fontsize=7.5, padding=1)
        for sp in ['top', 'right']:
            ax.spines[sp].set_visible(False)
        if ct == 'Neuron':
            ax.legend(fontsize=9, frameon=False, loc='upper left')
    fig.suptitle('Nucleus-only: same cells, two tests — Wilcoxon inflates the hit list',
                 fontsize=14, fontweight='bold')
    plt.tight_layout(rect=[0, 0, 1, 0.94])
    plt.savefig(out_png, dpi=170, bbox_inches='tight')
    plt.close(fig)


def render_pdf_png(pdf, out_png, dpi=200):
    """Rasterize the first page of a vector PDF figure to PNG (for embedding in
    the deck, which can't take PDF directly). Uses PyMuPDF."""
    import fitz
    doc = fitz.open(str(pdf))
    pix = doc[0].get_pixmap(matrix=fitz.Matrix(dpi / 72, dpi / 72))
    pix.save(str(out_png))
    doc.close()


def export_bar_genes_xlsx(v1_pb_csv, v1_wx_csv, nuc_pb_csv, nuc_wx_csv,
                          out_xlsx):
    """Source data for the two bar-chart slides: every gene counted in each bar,
    with its log2 fold change and adjusted p-value. One sheet per dataset, plus
    a Summary sheet whose counts reproduce the bar heights exactly.

    Significance rules match each figure:
      V1 Wilcoxon (PI_Master):  padj<1e-3 & |logfc|>0.5
      V1 pseudobulk:            the CSV 'significant' column
      nucleus Wilcoxon:         padj<1e-3 & |log2FoldChange|>0.5
      nucleus pseudobulk:       padj<0.05  & |log2FoldChange|>0.5
    """
    cols = ['cell_type', 'contrast', 'test', 'gene', 'log2FC', 'padj']

    def _v1():
        wx = pd.read_csv(v1_wx_csv)
        wx = wx[_wx_sig(wx)].copy()
        wx['contrast'] = wx['test'] + '_vs_' + wx['reference']
        keep = {f'{t}_vs_{r}' for t, r in V1_CONTRASTS}
        wx = wx[wx['contrast'].isin(keep)]
        wxo = wx.rename(columns={'subtype': 'cell_type', 'logfc': 'log2FC'})
        wxo['test'] = 'Cell-level Wilcoxon'
        wxo = wxo[cols]

        pb = pd.read_csv(v1_pb_csv)
        pb = pb[(pb['significant']) & (pb['comparison'].isin(keep))].copy()
        pbo = pb.rename(columns={'celltype': 'cell_type',
                                 'comparison': 'contrast',
                                 'log2FoldChange': 'log2FC'})
        pbo['test'] = 'Pseudobulk DESeq2'
        pbo = pbo[cols]
        return pd.concat([wxo, pbo], ignore_index=True)

    def _nuc():
        wx = pd.read_csv(nuc_wx_csv)
        wx = wx[(wx['padj'] < WX_PADJ) &
                (wx['log2FoldChange'].abs() > WX_LFC)].copy()
        pb = pd.read_csv(nuc_pb_csv)
        pb = pb[(pb['padj'] < PB_PADJ) &
                (pb['log2FoldChange'].abs() > PB_LFC)].copy()
        keep = {f'{t}_vs_{r}' for t, r in NUC_CONTRASTS}
        out = []
        for df, name in [(wx, 'Cell-level Wilcoxon'),
                         (pb, 'Pseudobulk DESeq2')]:
            d = df[df['comparison'].isin(keep)].rename(
                columns={'celltype': 'cell_type', 'comparison': 'contrast',
                         'log2FoldChange': 'log2FC'}).copy()
            d['test'] = name
            out.append(d[cols])
        return pd.concat(out, ignore_index=True)

    v1 = _v1().sort_values(['cell_type', 'contrast', 'test', 'padj'])
    nuc = _nuc().sort_values(['cell_type', 'contrast', 'test', 'padj'])

    def _summary(df, dataset):
        g = (df.groupby(['cell_type', 'contrast', 'test'])
             .size().reset_index(name='n_significant_genes'))
        g.insert(0, 'dataset', dataset)
        return g
    summary = pd.concat([_summary(v1, 'Whole-cell'),
                         _summary(nuc, 'Nucleus-only')], ignore_index=True)

    with pd.ExcelWriter(out_xlsx, engine='openpyxl') as xw:
        summary.to_excel(xw, sheet_name='Summary (bar heights)', index=False)
        v1.to_excel(xw, sheet_name='Whole-cell genes', index=False)
        nuc.to_excel(xw, sheet_name='Nucleus-only genes', index=False)
    print(f'[xlsx] {out_xlsx}  '
          f'(whole-cell {len(v1)} rows, nucleus {len(nuc)} rows)')


def compare_stats(csv, highlight=None):
    """Recompute the cell-typing-method comparison numbers straight from the
    compare CSV (cols: celltype, gene, logfc, padj, method, significant) so the
    deck slide never hardcodes them. Returns {celltype: {...}}."""
    df = pd.read_csv(csv)
    out = {}
    for ct in ['Neuron', 'Astrocyte']:
        sub = df[df['celltype'] == ct]
        m = set(sub[(sub['method'] == 'marker-typed') & (sub['significant'])]['gene'])
        l = set(sub[(sub['method'] == 'leiden-typed') & (sub['significant'])]['gene'])
        shared = m & l
        union = m | l
        jac = len(shared) / len(union) if union else float('nan')
        rec = {'marker_sig': len(m), 'leiden_sig': len(l),
               'shared': len(shared), 'jaccard': jac, 'hl': None}
        if highlight:
            hl = {}
            for meth in ['marker-typed', 'leiden-typed']:
                row = sub[(sub['method'] == meth) & (sub['gene'] == highlight)]
                if len(row):
                    rr = row.iloc[0]
                    hl[meth] = (float(rr['logfc']), float(rr['padj']),
                                bool(rr['significant']))
            rec['hl'] = hl
        out[ct] = rec
    return out


# ---- pptx helpers --------------------------------------------------------
def add_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def txt(slide, l, t, w, h, text, size=18, bold=False, color=NAVY,
        align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = 'Calibri'
    return tb


def bullets(slide, l, t, w, h, items, size=18, color=NAVY, space=10):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (lvl, text, *style) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(space)
        r = p.add_run()
        r.text = ('• ' if lvl == 0 else '– ') + text
        r.font.size = Pt(size - lvl * 2)
        r.font.color.rgb = style[0] if style else color
        r.font.name = 'Calibri'
        if len(style) > 1 and style[1]:
            r.font.bold = True
    return tb


def band(slide, prs, color=NAVY, h=1.05):
    s = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width,
                               Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def chip(slide, l, t, w, h, text, fill, tcolor=WHITE, size=13, bold=True):
    s = slide.shapes.add_shape(5, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = WHITE
    s.line.width = Pt(1)
    s.shadow.inherit = False
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = tcolor
    r.font.name = 'Calibri'
    return s


def arrow(slide, l, t, w, h=0.3, color=LGRAY):
    s = slide.shapes.add_shape(13, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def slide_header(slide, prs, title, sub=None, tag=None, tag_color=None):
    band(slide, prs)
    tw = 8.8 if tag else 12.3
    txt(slide, 0.5, 0.16, tw, 0.62, title, size=28, bold=True, color=WHITE)
    if sub:
        txt(slide, 0.5, 0.74, tw, 0.32, sub, size=13,
            color=RGBColor(0xCF, 0xDA, 0xE6))
    if tag:
        chip(slide, 9.5, 0.29, 3.4, 0.47, tag, tag_color or TEAL,
             tcolor=WHITE, size=11)


def style_table(table, header_size=13, body_size=12, bold_col0=True):
    for j, cell in enumerate(table.rows[0].cells):
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE
                r.font.bold = True
                r.font.size = Pt(header_size)
                r.font.name = 'Calibri'
    for i in range(1, len(table.rows)):
        for j in range(len(table.columns)):
            c = table.cell(i, j)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(body_size)
                    r.font.color.rgb = NAVY if j == 0 else GRAY
                    r.font.name = 'Calibri'
                    if j == 0 and bold_col0:
                        r.font.bold = True


def add_compare_slide(prs, png, csv, test, ref, highlight=None):
    """Cell-typing-method DE comparison slide. All numbers (sig counts, shared,
    Jaccard, highlight-gene status) are recomputed from `csv` — nothing
    hardcoded — so it works for any contrast."""
    st = compare_stats(csv, highlight)
    n, a = st['Neuron'], st['Astrocyte']
    s = add_blank(prs)
    slide_header(s, prs, 'Does the cell-typing method change the DE result?',
                 f'Same test (Wilcoxon), same contrast ({test} vs {ref}), '
                 'same thresholds — only HOW cells were typed differs',
                 tag='Whole-cell · marker vs leiden', tag_color=TEAL)
    s.shapes.add_picture(str(png), Inches(0.45), Inches(1.35),
                         height=Inches(5.75))
    chip(s, 7.35, 1.5, 5.5, 0.5, 'What it shows', PURPLE, size=14)
    items = [
        (0, 'Left column = marker-score typing (our DE pipeline); right = '
            'leiden clusters (the 10x run).', NAVY, True),
        (0, f"Neurons: {n['marker_sig']} vs {n['leiden_sig']} sig genes, "
            f"{n['shared']} shared (Jaccard {n['jaccard']:.2f}) — "
            'the DE answer barely moves.', BLUE, True),
        (0, f"Astrocytes: {a['marker_sig']} vs {a['leiden_sig']} sig, "
            f"{a['shared']} shared (Jaccard {a['jaccard']:.2f}) — "
            'noticeably more sensitive.', GREEN, True),
        (0, 'Why: neurons are abundant and cleanly separated, so both '
            'methods grab nearly the same cells.', NAVY),
        (0, 'Astrocytes are fewer and sit closer to other glia — the two '
            'methods disagree on the borderline cells.', NAVY),
    ]
    # if a highlight gene was requested, report its status both ways
    if highlight and n['hl']:
        def _fmt(d, ct):
            hl = st[ct]['hl']
            parts = []
            for meth, lab in [('marker-typed', 'marker'),
                              ('leiden-typed', 'leiden')]:
                if meth in hl:
                    lfc, padj, sig = hl[meth]
                    parts.append(f"{lab} log2FC={lfc:+.2f}, padj={padj:.1e} "
                                 f"({'SIG' if sig else 'n.s.'})")
            return '; '.join(parts)
        items.append(
            (0, f"{highlight} (Neuron): {_fmt(st, 'Neuron')}",
             RED, True))
    bullets(s, 7.4, 2.15, 5.55, 4.4, items, size=12.5)
    txt(s, 7.4, 6.5, 5.55, 0.85,
        'Takeaway: for big, distinct populations the DE is robust to the '
        'typing method; for sparser types it is not — define them carefully '
        'and state which method you used.',
        size=12, color=AMBER, bold=True)


# ---- build ---------------------------------------------------------------
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--pb-csv', required=True,
                    help='WholeCell_V1_PseudobulkDE_results CSV')
    ap.add_argument('--wx-csv', required=True,
                    help='PI_Master_NvsA_V1_DE Wilcoxon CSV')
    ap.add_argument('--out', required=True, help='output .pptx path')
    ap.add_argument('--compare', action='append', default=[],
                    help='repeatable cell-typing DE comparison spec, '
                         'PNG|CSV|TEST|REF[|HIGHLIGHT] '
                         '(from wholecell_V1_celltyping_DE_compare.py). '
                         'Pass once per contrast.')
    ap.add_argument('--nuc-pb-csv', default=None,
                    help='optional: nucleus-only Pseudobulk DE results CSV '
                         '(adds a nucleus version of the sig-bar slide)')
    ap.add_argument('--nuc-wx-csv', default=None,
                    help='optional: nucleus-only Wilcoxon DE results CSV')
    ap.add_argument('--umap-pdf', default=None,
                    help='optional: leiden UMAP PDF (WholeCell_V1_10x_UMAP_'
                         'leiden_slide) for the "how leiden worked" slide')
    ap.add_argument('--spatial-pdf', default=None,
                    help='optional: leiden spatial-cluster PDF '
                         '(WholeCell_V1_10x_spatial_clusters)')
    ap.add_argument('--wc-dotplot-pdf', default=None,
                    help='optional: whole-cell Neuron/Astrocyte iteration dot '
                         'plot PDF (CelltypeIteration_Dotplot..._WholeCell...)')
    ap.add_argument('--nuc-dotplot-pdf', default=None,
                    help='optional: nucleus-only Neuron/Astrocyte iteration dot '
                         'plot PDF (CelltypeIteration_Dotplot..._NucleusOnly...)')
    ap.add_argument('--date', default='2026-06-09')
    args = ap.parse_args()

    out = Path(args.out)
    D = args.date
    assets = out.parent / 'pipeline_deck_assets'
    assets.mkdir(exist_ok=True)

    sig_png = assets / 'sig_bar_v1.png'
    diag_png = assets / 'diagnostic_EtOHveh_neuron.png'
    print('[fig] significance bar chart...')
    fig_sig_bar(args.pb_csv, args.wx_csv, sig_png)
    print('[fig] underpowered-vs-null diagnostic...')
    fig_diagnostic(args.pb_csv, args.wx_csv, diag_png)

    leiden_umap_png = leiden_spatial_png = None
    if (args.umap_pdf and args.spatial_pdf and Path(args.umap_pdf).exists()
            and Path(args.spatial_pdf).exists()):
        leiden_umap_png = assets / 'leiden_umap.png'
        leiden_spatial_png = assets / 'leiden_spatial.png'
        print('[fig] rendering leiden UMAP + spatial maps from PDF...')
        render_pdf_png(args.umap_pdf, leiden_umap_png)
        render_pdf_png(args.spatial_pdf, leiden_spatial_png)

    wc_dot_png = nuc_dot_png = None
    if args.wc_dotplot_pdf and Path(args.wc_dotplot_pdf).exists():
        wc_dot_png = assets / 'dotplot_wholecell.png'
        print('[fig] rendering whole-cell iteration dot plot from PDF...')
        render_pdf_png(args.wc_dotplot_pdf, wc_dot_png, dpi=300)
    if args.nuc_dotplot_pdf and Path(args.nuc_dotplot_pdf).exists():
        nuc_dot_png = assets / 'dotplot_nucleus.png'
        print('[fig] rendering nucleus-only iteration dot plot from PDF...')
        render_pdf_png(args.nuc_dotplot_pdf, nuc_dot_png, dpi=300)

    nuc_sig_png = None
    if args.nuc_pb_csv and args.nuc_wx_csv:
        nuc_sig_png = assets / 'sig_bar_nucleus.png'
        print('[fig] nucleus-only significance bar chart...')
        fig_sig_bar_nucleus(args.nuc_pb_csv, args.nuc_wx_csv, nuc_sig_png)
        # source data for the two bar-chart slides (genes + log2FC + padj)
        bar_xlsx = out.parent / f'Xenium_BarChart_SignificantGenes_' \
                                f'WholeCell_Nucleus_{D}.xlsx'
        print('[xlsx] bar-chart significant-gene source data...')
        export_bar_genes_xlsx(args.pb_csv, args.wx_csv,
                              args.nuc_pb_csv, args.nuc_wx_csv, bar_xlsx)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: title ----
    s = add_blank(prs)
    bg = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False
    txt(s, 0.8, 2.1, 11.7, 1.4, 'Spatial transcriptomics with Claude',
        size=42, bold=True, color=WHITE)
    txt(s, 0.8, 3.5, 11.7, 0.9,
        'A practical 10x Xenium pipeline — and the caveats to watch for',
        size=22, color=RGBColor(0xCF, 0xDA, 0xE6))
    txt(s, 0.8, 6.45, 11.7, 0.5,
        f'{D}   ·   general overview', size=14, color=LGRAY)

    # ---- Slide 2: workflow (general, LARGE) ----
    s = add_blank(prs)
    slide_header(s, prs, 'The workflow: from your files to results')
    steps = [
        ('1\nHand over\nthe files', 'Xenium output folder\n(count matrix +\nmolecule positions)\n+ ROI CSVs from\nExplorer', PURPLE),
        ('2\nBuild the\nmatrix', 'cells × genes table;\nassign each cell to\nan ROI (point-in-\npolygon)', BLUE),
        ('3\nQC\nfilter', 'drop low-quality\ntranscripts, empty\ncells, rare genes', BLUE),
        ('4\nNormalize\n+ log', 'put cells on a\ncomparable scale', BLUE),
        ('5\nCell\ntyping', 'score marker genes\n→ Neuron, Astrocyte,\nOligo, Microglia', BLUE),
        ('6\nDifferential\nexpression', 'compare groups;\nvolcano plots +\nresult tables', GREEN),
    ]
    x = 0.55
    bw = 1.95
    gap = 0.12
    y = 1.9
    for i, (head, body, col) in enumerate(steps):
        chip(s, x, y, bw, 1.55, head, col, size=15)
        txt(s, x, y + 1.62, bw, 1.7, body, size=12.5, color=GRAY)
        if i < len(steps) - 1:
            arrow(s, x + bw + 0.002, y + 0.62, gap + 0.02)
        x += bw + gap
    txt(s, 0.55, 5.9, 12.2, 1.3,
        'Key idea: Xenium does not hand you "the data" as a table of genes per cell '
        'type. It gives molecule positions and a raw count matrix. Everything else — '
        'ROIs, cell types, expression per group — is computed downstream. You give the '
        'files and the questions; the pipeline builds the rest.',
        size=16, color=NAVY)

    # ---- Slide 2b: the iterations roadmap ----
    s = add_blank(prs)
    slide_header(s, prs, 'What this deck compares: three choices',
                 'Each comparison varies ONE of these and holds the rest fixed — '
                 'a tag in the top-right names which.')
    roadmap = [
        ('Segmentation', 'Whole-cell', TEAL, 'Nucleus-only', NUCP,
         'Which part of the cell we count transcripts in'),
        ('Cell typing\n(= the two pipelines)',
         'Marker scoring\n(V1 DE pipeline)', GREEN,
         'Leiden clusters\n(10x pipeline)', BLUE,
         'How each cell gets its identity'),
        ('DE test', 'Cell-level Wilcoxon', RED, 'Pseudobulk DESeq2', BLUE,
         'How we test treatment effects (pseudobulk = rigorous primary)'),
    ]
    y = 1.7
    for label, oa, ca, ob, cb, cap in roadmap:
        txt(s, 0.55, y + 0.02, 2.75, 0.95, label, size=15, bold=True, color=NAVY)
        chip(s, 3.5, y, 3.2, 0.78, oa, ca, size=13)
        txt(s, 6.78, y + 0.2, 0.55, 0.4, 'vs', size=14, bold=True, color=GRAY,
            align=PP_ALIGN.CENTER)
        chip(s, 7.45, y, 3.2, 0.78, ob, cb, size=13)
        txt(s, 10.78, y + 0.06, 2.4, 0.85, cap, size=10.5, color=GRAY)
        y += 1.45
    txt(s, 0.55, 6.45, 12.3, 0.9,
        'We do NOT run every combination. Each comparison slide changes just one '
        'of these three and keeps the rest fixed. The corner tag (e.g. '
        '"Whole-cell · Wilcoxon vs pseudobulk") names the iteration on screen. '
        'Pseudobulk is the rigorous primary test; the cell-typing comparisons that '
        'follow are Wilcoxon-based, so we trace that thread through to the end.',
        size=13, color=AMBER, bold=True)

    # ---- Slide 3: downstream QC, side-by-side (DE pipeline vs 10x workshop) ----
    s = add_blank(prs)
    slide_header(s, prs, 'Downstream QC: the DE pipeline vs the 10x workshop',
                 'Both run AFTER the instrument\'s QV>=20 transcript filter — '
                 'same raw data, different goals, different cutoffs.',
                 tag='Whole-cell', tag_color=TEAL)
    rows = [
        ('QC / processing step', 'V1 DE pipeline  (made the data)',
         '10x workshop pipeline  (just run)'),
        ('Min counts per cell', '>= 10', '>= 20'),
        ('Max counts per cell', 'none', '<= 98th percentile  (= 3,405)'),
        ('Min cells per gene', '>= 5', '>= 100'),
        ('Normalize', 'total -> 1e4, then log1p', 'total -> median, then log1p'),
        ('Highly-variable genes', 'not used — keep all genes',
         '2,000 (seurat_v3)'),
        ('Scaling', 'none', 'z-score, capped at 10'),
        ('Cell identity', 'score canonical markers -> label',
         'leiden clusters (unsupervised)'),
        ('Extra filters', 'pseudobulk gene >= 10 counts; replicate >= 25 cells',
         '—'),
        ('Result', '~214,744 cells · all 5,104 genes',
         '210,391 cells · 4,679 genes'),
    ]
    gt = s.shapes.add_table(len(rows), 3, Inches(0.5), Inches(1.55),
                            Inches(12.3), Inches(4.45)).table
    gt.columns[0].width = Inches(3.0)
    gt.columns[1].width = Inches(4.65)
    gt.columns[2].width = Inches(4.65)
    for i in range(len(rows)):
        for j in range(3):
            gt.cell(i, j).text = rows[i][j]
    style_table(gt, header_size=12.5, body_size=12)
    # colour the two pipeline columns to match the rest of the deck
    for i in range(1, len(rows)):
        for p in gt.cell(i, 1).text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = PURPLE
                r.font.bold = True
        for p in gt.cell(i, 2).text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = GREEN
                r.font.bold = True
    txt(s, 0.5, 6.2, 12.3, 1.2,
        'Different goals, different cutoffs. DE counts molecules per mouse, so QC stays '
        'permissive — keep every real cell and gene. Clustering hunts for structure, so it '
        'trims harder, keeps only the 2,000 most variable genes, and scales. Neither set is '
        '"the" 10x number: 10x says set cell/gene cutoffs from each dataset\'s own distributions.',
        size=13, color=NAVY)

    # ---- Slide 4: filter by cluster ----
    s = add_blank(prs)
    slide_header(s, prs, 'Optional: pull just one cluster / cell population',
                 'The cluster labels are already in the output files — no extra export needed')
    chip(s, 0.55, 1.7, 5.9, 0.55, 'In Xenium Explorer (to look)', PURPLE, size=14)
    bullets(s, 0.6, 2.45, 6.0, 3.6, [
        (0, 'Open the Cell groups panel (right side).', NAVY),
        (0, 'Click the cluster you want → "Select cells in group".', NAVY),
        (0, 'Useful for viewing, but the selection CSV is just outlines / '
            'coordinates — no counts.', GRAY),
    ], size=15)
    chip(s, 6.85, 1.7, 5.95, 0.55, 'Straight from the output (recommended)', BLUE, size=14)
    bullets(s, 6.9, 2.45, 5.95, 3.8, [
        (0, 'Cluster assignments already ship in each slide:', NAVY, True),
        (1, 'analysis/clustering/', GRAY),
        (1, 'gene_expression_graphclust/clusters.csv', GRAY),
        (0, 'Two columns — Barcode (cell ID) + Cluster. Filter Cluster == 1.', NAVY),
        (0, '10 clusterings ship: graphclust (Explorer default) + kmeans 2–10 '
            '— just say which.', GRAY),
    ], size=15)
    txt(s, 0.55, 6.15, 12.2, 1.1,
        'So Claude can pull cluster 1 directly from the output folder — no Explorer '
        'export needed for cluster identity. Just specify the clustering (almost always '
        'graphclust). Expression always comes from cell_feature_matrix.h5.',
        size=15, color=AMBER, bold=True)

    # ---- Slide 4d: two ways to cell type ----
    s = add_blank(prs)
    slide_header(s, prs, 'Two ways to assign cell types',
                 'Both label the same cells — one leans on an annotated '
                 'reference, the other on known marker genes',
                 tag='Choice: cell typing', tag_color=SLATE)
    chip(s, 0.55, 1.55, 5.95, 0.55, '1.  Reference mapping (label transfer)',
         PURPLE, size=14)
    bullets(s, 0.6, 2.3, 6.0, 4.2, [
        (0, 'Start from an ALREADY-annotated reference dataset (an atlas).',
            NAVY, True),
        (0, 'Project your cells into a shared space; each cell inherits the '
            'label of the reference cells it most resembles.', NAVY),
        (0, 'Tools: Azimuth, Seurat TransferData, scANVI / scArches, '
            'scanpy ingest.', GRAY),
        (0, 'Pros: automated, consistent with the reference taxonomy, good for '
            'fine subtypes.', NAVY),
        (0, 'Needs: a high-quality reference matching your tissue, species '
            'and platform.', GRAY),
    ], size=14)
    chip(s, 6.85, 1.55, 5.95, 0.55, '2.  Marker gene annotation', BLUE, size=14)
    bullets(s, 6.9, 2.3, 5.95, 4.2, [
        (0, 'Use known marker genes to decide identity — NO reference needed.',
            NAVY, True),
        (0, 'Cluster-then-label: cluster cells (leiden), read each cluster\'s '
            'top markers, name it.', NAVY),
        (1, '← the 10x / leiden run did this', BLUE, True),
        (0, 'Per-cell scoring: score every cell against marker sets, take the '
            'highest (score_genes → argmax).', NAVY),
        (1, '← the V1 DE pipeline did this', GREEN, True),
        (0, 'Pros: transparent, works without a matching reference. '
            'Needs: curated marker lists + human judgement.', GRAY),
    ], size=14)
    txt(s, 0.55, 6.45, 12.3, 0.95,
        'Both pipelines here used marker gene annotation — neither did reference '
        'mapping. V1 scored each cell on canonical markers; the 10x run clustered '
        'with leiden first, then labelled each cluster. Same family, two styles — '
        'which is exactly why their DE results line up so closely.',
        size=13, color=AMBER, bold=True)

    # ---- Slide 4e: how the leiden clustering worked ----
    s = add_blank(prs)
    slide_header(s, prs, 'How the leiden clustering worked',
                 'From the expression matrix to 19 spatial clusters — the '
                 'unsupervised recipe behind the 10x run',
                 tag='Whole-cell · leiden', tag_color=TEAL)
    chip(s, 0.5, 1.4, 4.35, 0.5, 'The recipe (no labels used)', PURPLE, size=13)
    bullets(s, 0.55, 2.0, 4.4, 3.3, [
        (0, 'Normalize + log; keep the 2,000 most variable genes.', NAVY),
        (0, 'PCA → top 30 components (denoise).', NAVY),
        (0, 'Build a k-nearest-neighbour graph — link each cell to the cells '
            'most similar to it.', NAVY),
        (0, 'Leiden finds densely-connected COMMUNITIES in that graph → '
            '19 clusters (resolution 0.5).', NAVY, True),
        (0, 'UMAP = a 2-D picture of the same graph, for viewing only — it is '
            'NOT what does the clustering.', GRAY),
        (0, 'Map each cluster back onto the cell\'s (x, y) → the tissue view.',
            NAVY),
    ], size=12.5)
    txt(s, 0.55, 5.55, 4.4, 1.6,
        'Leiden returns 19 anonymous clusters — colours, not names. Turning '
        'those colours into cell types is the SEPARATE step on the next slide. '
        'Raise the resolution → more clusters; same recipe.',
        size=12, color=AMBER, bold=True)
    txt(s, 5.05, 1.2, 7.95, 0.3,
        'UMAP — every cell placed by expression; colour = leiden cluster (left), '
        'slide (right)', size=10.5, bold=True, color=NAVY)
    if leiden_umap_png and leiden_umap_png.exists():
        s.shapes.add_picture(str(leiden_umap_png), Inches(5.05), Inches(1.5),
                             width=Inches(7.95))
    txt(s, 5.05, 3.95, 7.95, 0.3,
        'The same clusters mapped back onto each tissue section',
        size=10.5, bold=True, color=NAVY)
    if leiden_spatial_png and leiden_spatial_png.exists():
        s.shapes.add_picture(str(leiden_spatial_png), Inches(6.35), Inches(4.25),
                             width=Inches(4.9))

    # ---- Slide 4c: cell typing (cluster -> cell-type name) ----
    s = add_blank(prs)
    slide_header(s, prs, 'Cell typing: putting names on the clusters',
                 'Clustering gives anonymous groups; the cell-type label is a '
                 'separate, human-checked step')
    chip(s, 0.55, 1.5, 5.6, 0.5, 'From clusters to cell types', PURPLE, size=14)
    bullets(s, 0.6, 2.12, 5.7, 2.3, [
        (0, 'Clustering only gives numbered groups (0, 1, 2 …) — never names.',
            NAVY, True),
        (0, 'For each cluster, list the genes it turns UP vs all other cells.',
            NAVY),
        (0, 'Match that marker signature to known cell-type markers (right).',
            NAVY),
        (0, 'Assign the label; merge clusters that share an identity, split '
            'one that holds two.', NAVY),
        (0, 'Shortcut: score every cell for each marker set, take the highest.',
            GRAY),
    ], size=13.5)
    chip(s, 0.55, 4.55, 5.6, 0.5, 'Where the marker lists come from', BLUE,
         size=14)
    bullets(s, 0.6, 5.15, 5.7, 1.5, [
        (0, 'Curated databases: PanglaoDB, CellMarker 2.0', GRAY),
        (0, 'Brain-specific: Allen Brain Atlas / ABC Atlas', GRAY),
        (0, 'Automated: Azimuth label-transfer vs an annotated reference', GRAY),
    ], size=12.5)
    chip(s, 6.45, 1.5, 6.4, 0.5, 'Brain marker cheat-sheet (mouse)', NAVY,
         size=14)
    mk = [
        ('Cell type', 'Canonical markers'),
        ('Excitatory neuron', 'Snap25, Slc17a7, Rbfox3'),
        ('Inhibitory neuron', 'Gad1, Gad2, Slc32a1'),
        ('Astrocyte', 'Aqp4, Gfap, Slc1a3, Gja1'),
        ('Oligodendrocyte', 'Mbp, Mog, Mag, Cnp'),
        ('OPC (oligo precursor)', 'Pdgfra, Olig2, Cspg4'),
        ('Microglia', 'Csf1r, Cx3cr1, P2ry12'),
        ('Endothelial', 'Cldn5, Slc2a1, Pecam1'),
        ('Pericyte / mural', 'Pdgfrb, Notch3, Vtn'),
    ]
    mt = s.shapes.add_table(len(mk), 2, Inches(6.45), Inches(2.15),
                            Inches(6.4), Inches(3.95)).table
    mt.columns[0].width = Inches(2.7)
    mt.columns[1].width = Inches(3.7)
    for i in range(len(mk)):
        for j in range(2):
            mt.cell(i, j).text = mk[i][j]
    style_table(mt, header_size=12.5, body_size=11.5)
    txt(s, 0.55, 6.75, 12.3, 0.7,
        'These eight are exactly what tagged our 19 clusters (Cnp/Mag/Mog → '
        'oligo, Gad1/Gad2 → inhibitory, Slc1a3/Gja1 → astrocyte …). The names '
        'are interpretation, not output — always sanity-check against an atlas '
        'or an expert, and confirm a cluster isn\'t just one mouse or one slide.',
        size=12.5, color=AMBER, bold=True)

    # ---- Slide 4b: spatial neighborhood / niche analysis ----
    s = add_blank(prs)
    slide_header(s, prs, 'Spatial neighborhood / niche analysis',
                 'Which cells sit near a given cell type — and what recurrent '
                 'microenvironments exist',
                 tag='Whole-cell', tag_color=TEAL)
    chip(s, 0.55, 1.6, 5.9, 0.55, 'Questions you can ask', PURPLE, size=14)
    bullets(s, 0.6, 2.3, 6.0, 4.2, [
        (0, 'Is cell type A enriched right next to cell type B?', NAVY),
        (0, 'What is the cell-type mix surrounding each astrocyte (or neuron)?', NAVY),
        (0, 'Do neurons NEAR astrocytes express different genes than neurons far away?', NAVY),
        (0, 'What recurrent tissue "niches" exist across the section?', NAVY),
    ], size=15)
    chip(s, 6.85, 1.6, 5.95, 0.55, 'How it is done (the recipe)', BLUE, size=14)
    bullets(s, 6.9, 2.3, 5.95, 4.2, [
        (0, 'Inputs we ALREADY have: cell centroids (x,y) + cell-type labels.', NAVY, True),
        (0, 'Build a spatial neighbor graph:', NAVY),
        (1, 'k-nearest neighbours (k ≈ 6–20), or', GRAY),
        (1, 'all cells within a fixed radius (≈ 30–50 µm)', GRAY),
        (0, 'Enrichment: permutation test — observed vs random adjacency.', NAVY),
        (0, 'Niches: summarise each cell\'s neighbour composition → cluster.', NAVY),
    ], size=15)
    txt(s, 0.55, 6.35, 12.2, 1.0,
        '10x ships no niche tool onboard — use Squidpy (Python; runs on our AnnData '
        'object) or Seurat BuildNicheAssay (R). Caveat: build the graph PER SECTION '
        '(never link cells across slides), and coordinates are in microns.',
        size=14, color=AMBER, bold=True)

    # ---- Slide 5: Wilcoxon vs pseudobulk (why each) ----
    s = add_blank(prs)
    slide_header(s, prs, 'Two ways to test differential expression',
                 'Pick the test that matches your replicate structure',
                 tag='Choice: DE test', tag_color=SLATE)
    chip(s, 0.55, 1.55, 5.95, 0.55, 'Cell-level Wilcoxon', RED, size=15)
    bullets(s, 0.6, 2.25, 6.0, 4.6, [
        (0, 'Treats every CELL as a sample.', NAVY, True),
        (0, 'Fast; good for marker discovery and exploration.', NAVY),
        (0, 'Reasonable when you truly have n=1 (e.g. one punch).', NAVY),
        (0, 'BUT with 2–3 mice it is pseudoreplication: p-values track '
            'cell COUNT, not biology → inflated hit lists.', RED, True),
        (0, 'Reviewers discount these p-values.', GRAY),
    ], size=15)
    chip(s, 6.85, 1.55, 5.95, 0.55, 'Pseudobulk DESeq2', BLUE, size=15)
    bullets(s, 6.9, 2.25, 5.95, 4.6, [
        (0, 'Sums counts per MOUSE, then tests across animals.', NAVY, True),
        (0, 'Mouse = the real experimental unit.', NAVY),
        (0, 'Rigorous and reviewer-preferred when n≥2–3.', BLUE, True),
        (0, 'Conservative: needs the effect to reproduce across animals.', NAVY),
        (0, 'Use this as the PRIMARY analysis.', BLUE, True),
    ], size=15)

    # ---- Slide 6: the diagnostic (null != no effect) ----
    s = add_blank(prs)
    slide_header(s, prs, 'A "no significant DE" result is not always "no effect"',
                 'Worked example: EtOH_veh vs H2O_veh, neurons (n=2 vs 2)',
                 tag='Whole-cell · pseudobulk', tag_color=TEAL)
    s.shapes.add_picture(str(diag_png), Inches(0.7), Inches(1.5),
                         width=Inches(12.0))
    txt(s, 0.55, 6.55, 12.2, 0.85,
        'Hundreds of large fold changes, all pointing the same way as the Wilcoxon '
        'result, capped just under significance by sample size. That is an '
        'underpowered effect — the fix is more animals, not a looser test.',
        size=14, color=NAVY)

    # ---- Slide 7: the bar chart ----
    s = add_blank(prs)
    slide_header(s, prs, 'Same data, two tests — the size of the inflation',
                 'Significant genes per contrast (|log2FC|>0.5), V1 whole-cell',
                 tag='Whole-cell · Wilcoxon vs pseudobulk', tag_color=TEAL)
    s.shapes.add_picture(str(sig_png), Inches(0.6), Inches(1.45),
                         width=Inches(12.1))
    txt(s, 0.55, 6.35, 12.2, 1.0,
        'Wilcoxon (red) returns far more "hits" than pseudobulk (blue) — even though '
        'pseudobulk uses the looser p-threshold here. The gap is pseudoreplication. '
        'The contrasts that survive pseudobulk (e.g. the MAT2A and EtOH_MCT1i comparisons) '
        'are the ones with the strongest, most reproducible effects.',
        size=13.5, color=NAVY)

    # ---- Slide 7n: nucleus-only version of the bar chart ----
    if nuc_sig_png is not None:
        s = add_blank(prs)
        slide_header(s, prs, 'Nucleus-only: same data, two tests',
                     'Significant genes per contrast (|log2FC|>0.5), '
                     'nucleus-only segmentation (all vs H2O_veh)',
                     tag='Nucleus-only · Wilcoxon vs pseudobulk',
                     tag_color=NUCP)
        s.shapes.add_picture(str(nuc_sig_png), Inches(0.6), Inches(1.45),
                             width=Inches(12.1))
        txt(s, 0.55, 6.35, 12.2, 1.0,
            'Same pattern holds on the nucleus-only data: cell-level Wilcoxon (red) '
            'returns far more "hits" than pseudobulk (blue) under the looser p-cut. '
            'The inflation is pseudoreplication, not a property of the segmentation — '
            'pseudobulk stays the reviewer-defensible primary analysis.',
            size=13.5, color=NAVY)

    # ---- Slide 7d/7e: Neuron vs Astrocyte dot plots, one per segmentation iteration ----
    def add_dotplot_slide(png, seg_label, tag, tag_color, lead):
        s = add_blank(prs)
        slide_header(
            s, prs, 'Curated gene panel — Neuron vs Astrocyte dot plot',
            f'{seg_label}: dot size = % cells positive · colour = Wilcoxon '
            'log2FC (clipped ±1) · all contrasts vs H2O_veh',
            tag=tag, tag_color=tag_color)
        # tall portrait figure on the left
        s.shapes.add_picture(str(png), Inches(0.4), Inches(1.3),
                             height=Inches(5.95))
        chip(s, 4.55, 1.45, 8.3, 0.5, 'How to read it', PURPLE, size=14)
        bullets(s, 4.6, 2.1, 8.25, 3.6, [
            (0, 'Two panels = the two cell types; four columns = the four '
                'treatment contrasts (all vs H2O_veh).', NAVY, True),
            (0, 'Dot SIZE = % of cells in the test group expressing the gene '
                '(detection frequency).', NAVY),
            (0, 'Dot COLOUR = Wilcoxon log2 fold change vs H2O_veh — red up, '
                'blue down, clipped at ±1.', NAVY),
            (0, 'Rows are grouped: neuronal markers, astrocyte markers, '
                'activity/receptor genes, acetyl-CoA/chromatin enzymes, and '
                'other alcohol-related genes.', GRAY),
            (0, 'Sanity check: neuronal markers light up (big dots) in the '
                'Neuron panel, astrocyte markers in the Astrocyte panel.', GREEN,
                True),
        ], size=13)
        txt(s, 4.6, 5.95, 8.25, 1.35, lead, size=13, color=AMBER, bold=True)
        return s

    if wc_dot_png is not None:
        add_dotplot_slide(
            wc_dot_png, 'WHOLE-CELL iteration',
            'Whole-cell · Neuron vs Astrocyte', TEAL,
            'Whole-cell counts pick up the strong acute-ethanol signal — e.g. '
            'Grin2d is up in neurons under EtOH. Compare this panel head-to-head '
            'with the nucleus-only version on the next slide.')
    if nuc_dot_png is not None:
        add_dotplot_slide(
            nuc_dot_png, 'NUCLEUS-ONLY iteration',
            'Nucleus-only · Neuron vs Astrocyte', NUCP,
            'Same genes, same axes, nucleus-only segmentation (the iteration '
            'with a single DE version — Wilcoxon). Dots are generally smaller '
            '(fewer transcripts per nucleus) and fold changes are muted versus '
            'whole-cell — the cost of counting only the nucleus.')

    # ---- Slide 7b: does the cell-typing method change DE? (one per --compare) ----
    for spec in args.compare:
        parts = spec.split('|')
        if len(parts) < 4:
            print(f'[warn] skipping malformed --compare spec: {spec}')
            continue
        png, csv, test, ref = parts[0], parts[1], parts[2], parts[3]
        highlight = parts[4] if len(parts) > 4 and parts[4] else None
        if not Path(png).exists() or not Path(csv).exists():
            print(f'[warn] skipping --compare, missing file(s): {png} / {csv}')
            continue
        print(f'[slide] compare {test} vs {ref}'
              + (f' (highlight {highlight})' if highlight else ''))
        add_compare_slide(prs, png, csv, test, ref, highlight)

    # ---- Slide 8: takeaways ----
    s = add_blank(prs)
    slide_header(s, prs, 'Takeaways')
    bullets(s, 0.6, 1.55, 12.2, 4.8, [
        (0, 'You provide the files + the question; the pipeline builds matrix → ROIs '
            '→ QC → cell types → DE.', NAVY, True),
        (0, 'The machine filters transcript quality (QV>=20) for you; you still must QC '
            'cells and genes — set cutoffs from the data and keep them identical across samples.', NAVY),
        (0, 'Match the DE test to your design: pseudobulk when you have biological '
            'replicates; Wilcoxon only for exploration / true n=1.', BLUE, True),
        (0, 'A null pseudobulk result with large, consistent fold changes = underpowered '
            '→ add Ns. Fold changes hugging zero = likely truly no effect.', NAVY),
        (0, 'Watch for confounds: if a treatment lives on only one slide/batch, its "hits" '
            'may be batch, not biology.', AMBER, True),
        (0, 'Every figure ships with editable text + a source-data CSV; scripts are version-'
            'controlled.', GRAY),
    ], size=17)

    prs.save(str(out))
    print(f'[saved] {out}')
    print(f'[assets] {sig_png.name}, {diag_png.name}')


if __name__ == '__main__':
    main()
